"""Concrete content extractor implementations for PDF, PPTX, and DOCX files.
 
Handles the following visual content types per file format:
- PDF: embedded images (passed as bytes to Claude vision), scanned pages (zero-text pages sent as images)
- PPTX: native charts (data extracted from XML), inserted picture shapes (bytes to Claude vision)
- DOCX: embedded images (bytes to Claude vision)
 
Mathematical symbols and special unicode characters are normalised before chunking.
"""
 
from __future__ import annotations
 
import io
import unicodedata
 
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
 
from BE.client.databricks_llm_client import DatabricksLLMClient
from BE.models.enums import FileType
from BE.patterns.extraction_strategy import ContentExtractor, ExtractedContent
from BE.utils.logger import get_logger
 
logger = get_logger(__name__)
 
_IMAGE_EXTRACTION_PROMPT = (
    "You are a financial document analyst. Describe the data in this image/chart in plain text, "
    "including all numbers, labels, axes, and trends. Output structured text only."
)
 
_UNICODE_NORMALIZE_FORM = "NFKC"
 
 
# ---------------------------------------------------------------------------
# PDF Extractor
# ---------------------------------------------------------------------------
 
 
class PdfExtractor(ContentExtractor):
    """Extracts content from PDF files using pdfplumber."""
 
    def __init__(self, llm_client: DatabricksLLMClient | None = None) -> None:
        self._llm_client = llm_client
 
    async def extract(self, file_content: bytes, file_name: str) -> list[ExtractedContent]:
        results: list[ExtractedContent] = []
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                # --- Text ---
                raw_text = page.extract_text() or ""
                text = _normalise_text(raw_text)
                has_text = bool(text.strip())
                if has_text:
                    results.append(
                        ExtractedContent(
                            chunk_text=text.strip(),
                            page_number=page_num,
                            chunk_type="text",
                            metadata={"source": file_name},
                        )
                    )
 
                # --- Tables ---
                tables = page.extract_tables() or []
                for idx, table in enumerate(tables):
                    table_text = _format_table(table)
                    if table_text.strip():
                        results.append(
                            ExtractedContent(
                                chunk_text=table_text,
                                page_number=page_num,
                                chunk_type="table",
                                metadata={"source": file_name, "table_index": idx},
                            )
                        )
 
                if not self._llm_client:
                    continue
 
                # --- Scanned page: no text but the page itself is an image ---
                # Render the full page to PNG and send to Claude vision
                if not has_text:
                    try:
                        page_image_bytes = _render_pdf_page_to_png(page)
                        if page_image_bytes:
                            description = await self._llm_client.invoke_sonnet_with_image(
                                system_prompt=_IMAGE_EXTRACTION_PROMPT,
                                image_bytes=page_image_bytes,
                                media_type="image/png",
                                text_prompt=(
                                    f"This is a scanned page {page_num} from '{file_name}'. "
                                    "Extract all text, numbers, and data visible on this page."
                                ),
                            )
                            content_text = _extract_llm_text(description)
                            if content_text.strip():
                                results.append(
                                    ExtractedContent(
                                        chunk_text=content_text,
                                        page_number=page_num,
                                        chunk_type="text",
                                        metadata={"source": file_name, "scanned": True},
                                    )
                                )
                    except Exception:
                        logger.warning(
                            "Failed to extract scanned page %d of %s", page_num, file_name, exc_info=True
                        )
                    continue  # scanned page handled — skip embedded image loop
 
                # --- Embedded images / charts on a text page ---
                # Extract actual image bytes and send to Claude vision
                for img_meta in page.images:
                    try:
                        image_bytes = _extract_pdf_image_bytes(page, img_meta)
                        if not image_bytes:
                            continue
                        description = await self._llm_client.invoke_sonnet_with_image(
                            system_prompt=_IMAGE_EXTRACTION_PROMPT,
                            image_bytes=image_bytes,
                            media_type="image/png",
                            text_prompt=f"Describe this chart or image from page {page_num} of '{file_name}'.",
                        )
                        content_text = _extract_llm_text(description)
                        if content_text.strip():
                            results.append(
                                ExtractedContent(
                                    chunk_text=content_text,
                                    page_number=page_num,
                                    chunk_type="chart",
                                    metadata={"source": file_name},
                                )
                            )
                    except Exception:
                        logger.warning(
                            "Failed to extract image on page %d of %s", page_num, file_name, exc_info=True
                        )
 
        return results
 
 
# ---------------------------------------------------------------------------
# PPTX Extractor
# ---------------------------------------------------------------------------
 
 
class PptxExtractor(ContentExtractor):
    """Extracts content from PPTX files using python-pptx."""
 
    def __init__(self, llm_client: DatabricksLLMClient | None = None) -> None:
        self._llm_client = llm_client
 
    async def extract(self, file_content: bytes, file_name: str) -> list[ExtractedContent]:
        results: list[ExtractedContent] = []
        prs = Presentation(io.BytesIO(file_content))
 
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
 
            for shape in slide.shapes:
                # --- Tables ---
                if shape.has_table:
                    table_text = _format_pptx_table(shape.table)
                    if table_text.strip():
                        results.append(
                            ExtractedContent(
                                chunk_text=table_text,
                                page_number=slide_num,
                                chunk_type="table",
                                metadata={"source": file_name},
                            )
                        )
                    continue
 
                # --- Native PowerPoint charts ---
                # Extract series/category/value data directly from chart XML — no LLM needed
                if shape.has_chart:
                    chart_text = _extract_pptx_chart_data(shape.chart)
                    if chart_text.strip():
                        results.append(
                            ExtractedContent(
                                chunk_text=chart_text,
                                page_number=slide_num,
                                chunk_type="chart",
                                metadata={"source": file_name},
                            )
                        )
                    continue
 
                # --- Inserted picture shapes (screenshots, photos, image-based graphs) ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self._llm_client:
                    try:
                        image_bytes = shape.image.blob
                        media_type = _pptx_image_media_type(shape.image.content_type)
                        description = await self._llm_client.invoke_sonnet_with_image(
                            system_prompt=_IMAGE_EXTRACTION_PROMPT,
                            image_bytes=image_bytes,
                            media_type=media_type,
                            text_prompt=f"Describe this image from slide {slide_num} of '{file_name}'.",
                        )
                        content_text = _extract_llm_text(description)
                        if content_text.strip():
                            results.append(
                                ExtractedContent(
                                    chunk_text=content_text,
                                    page_number=slide_num,
                                    chunk_type="chart",
                                    metadata={"source": file_name},
                                )
                            )
                    except Exception:
                        logger.warning(
                            "Failed to extract picture shape on slide %d of %s",
                            slide_num, file_name, exc_info=True,
                        )
                    continue
 
                # --- Text frames ---
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = _normalise_text(paragraph.text).strip()
                        if text:
                            slide_texts.append(text)
 
            if slide_texts:
                results.append(
                    ExtractedContent(
                        chunk_text="\n".join(slide_texts),
                        page_number=slide_num,
                        chunk_type="text",
                        metadata={"source": file_name},
                    )
                )
 
            # --- Speaker notes ---
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = _normalise_text(slide.notes_slide.notes_text_frame.text).strip()
                if notes_text:
                    results.append(
                        ExtractedContent(
                            chunk_text=notes_text,
                            page_number=slide_num,
                            chunk_type="text",
                            metadata={"source": file_name, "is_speaker_notes": True},
                        )
                    )
 
        return results
 
 
# ---------------------------------------------------------------------------
# DOCX Extractor
# ---------------------------------------------------------------------------
 
 
class DocxExtractor(ContentExtractor):
    """Extracts content from DOCX files using python-docx."""
 
    def __init__(self, llm_client: DatabricksLLMClient | None = None) -> None:
        self._llm_client = llm_client
 
    async def extract(self, file_content: bytes, file_name: str) -> list[ExtractedContent]:
        results: list[ExtractedContent] = []
        doc = DocxDocument(io.BytesIO(file_content))
 
        current_section: str | None = None
        section_texts: list[str] = []
 
        for paragraph in doc.paragraphs:
            if paragraph.style and paragraph.style.name and paragraph.style.name.startswith("Heading"):
                if section_texts:
                    results.append(
                        ExtractedContent(
                            chunk_text="\n".join(section_texts),
                            section_name=current_section,
                            chunk_type="text",
                            metadata={"source": file_name},
                        )
                    )
                    section_texts = []
                current_section = paragraph.text.strip() or None
                continue
 
            text = _normalise_text(paragraph.text).strip()
            if text:
                section_texts.append(text)
 
        if section_texts:
            results.append(
                ExtractedContent(
                    chunk_text="\n".join(section_texts),
                    section_name=current_section,
                    chunk_type="text",
                    metadata={"source": file_name},
                )
            )
 
        # --- Tables ---
        for idx, table in enumerate(doc.tables):
            table_text = _format_docx_table(table)
            if table_text.strip():
                results.append(
                    ExtractedContent(
                        chunk_text=table_text,
                        section_name=current_section,
                        chunk_type="table",
                        metadata={"source": file_name, "table_index": idx},
                    )
                )
 
        # --- Embedded images — extract actual bytes and send to Claude vision ---
        if self._llm_client:
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_bytes = rel.target_part.blob
                        media_type = rel.target_part.content_type or "image/png"
                        description = await self._llm_client.invoke_sonnet_with_image(
                            system_prompt=_IMAGE_EXTRACTION_PROMPT,
                            image_bytes=image_bytes,
                            media_type=media_type,
                            text_prompt=f"Describe this embedded image or chart from '{file_name}'.",
                        )
                        content_text = _extract_llm_text(description)
                        if content_text.strip():
                            results.append(
                                ExtractedContent(
                                    chunk_text=content_text,
                                    chunk_type="chart",
                                    metadata={"source": file_name},
                                )
                            )
                    except Exception:
                        logger.warning("Failed to extract image from %s", file_name, exc_info=True)
 
        return results
 
 
# ---------------------------------------------------------------------------
# Factory
# ---------------------------------------------------------------------------
 
 
class ContentExtractorFactory:
    """Returns the appropriate ContentExtractor for a given FileType."""
 
    @staticmethod
    def get_extractor(file_type: FileType, llm_client: DatabricksLLMClient | None = None) -> ContentExtractor:
        extractors: dict[FileType, type[ContentExtractor]] = {
            FileType.pdf: PdfExtractor,
            FileType.pptx: PptxExtractor,
            FileType.docx: DocxExtractor,
        }
        extractor_cls = extractors.get(file_type)
        if extractor_cls is None:
            raise ValueError(f"Unsupported file type: {file_type}")
        return extractor_cls(llm_client=llm_client)
 
 
# ---------------------------------------------------------------------------
# Helpers — formatting
# ---------------------------------------------------------------------------
 
 
def _format_table(table: list[list[str | None]]) -> str:
    """Format a pdfplumber table (list of rows) into readable text."""
    rows = []
    for row in table:
        cells = [str(cell) if cell is not None else "" for cell in row]
        rows.append(" | ".join(cells))
    return "\n".join(rows)
 
 
def _format_pptx_table(table) -> str:
    """Format a python-pptx Table object into readable text."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    return "\n".join(rows)
 
 
def _format_docx_table(table) -> str:
    """Format a python-docx Table object into readable text."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    return "\n".join(rows)
 
 
def _extract_llm_text(response: dict) -> str:
    """Extract text content from an LLM response dict."""
    try:
        return response["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError):
        return str(response)
 
 
# ---------------------------------------------------------------------------
# Helpers — text normalisation (Fix 5: math symbols / special characters)
# ---------------------------------------------------------------------------
 
 
def _normalise_text(text: str) -> str:
    """Normalise unicode text using NFKC form.
 
    NFKC converts compatibility characters to their canonical equivalents:
    - Mathematical symbols: ≥ ≤ ∑ × ÷ kept as readable unicode
    - Ligatures: fi fl restored from ligature characters
    - Superscripts/subscripts: ² → 2
    - Fullwidth characters: ａ → a
    - Fraction characters: ½ → 1/2
    """
    return unicodedata.normalize(_UNICODE_NORMALIZE_FORM, text)
 
 
# ---------------------------------------------------------------------------
# Helpers — PDF image extraction (Fix 2: actual image bytes)
# ---------------------------------------------------------------------------
 
 
def _extract_pdf_image_bytes(page, img_meta: dict) -> bytes | None:
    """Crop the image bounding box from a pdfplumber page and return PNG bytes."""
    try:
        x0, y0, x1, y1 = img_meta["x0"], img_meta["y0"], img_meta["x1"], img_meta["y1"]
        cropped = page.crop((x0, y0, x1, y1))
        pil_image = cropped.to_image(resolution=150).original
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return None
 
 
def _render_pdf_page_to_png(page, resolution: int = 150) -> bytes | None:
    """Render a full PDF page to PNG bytes for scanned page handling."""
    try:
        pil_image = page.to_image(resolution=resolution).original
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return None
 
 
# ---------------------------------------------------------------------------
# Helpers — PPTX chart data extraction (Fix 1: actual chart data from XML)
# ---------------------------------------------------------------------------
 
 
def _extract_pptx_chart_data(chart) -> str:
    """Extract series names, categories, and values from a native PowerPoint chart.
 
    Reads the underlying chart XML via python-pptx's chart data model and returns
    a structured plain-text representation — no LLM call required.
    """
    lines: list[str] = []
 
    try:
        chart_title = ""
        if chart.has_title and chart.chart_title.has_text_frame:
            chart_title = chart.chart_title.text_frame.text.strip()
        if chart_title:
            lines.append(f"Chart: {chart_title}")
 
        for plot in chart.plots:
            categories: list[str] = []
            try:
                cat_seq = getattr(plot, "categories", None)
                if cat_seq is not None:
                    categories = [str(c) for c in cat_seq]
            except Exception:
                pass
 
            for series in plot.series:
                series_name = getattr(series, "name", None) or "Series"
                values: list = []
                try:
                    values = [v for v in series.values if v is not None]
                except Exception:
                    pass
 
                if categories and values:
                    lines.append(f"Series: {series_name}")
                    for cat, val in zip(categories, values):
                        lines.append(f"  {cat}: {val}")
                elif values:
                    lines.append(f"Series: {series_name} — values: {', '.join(str(v) for v in values)}")
 
    except Exception:
        logger.warning("Failed to extract chart data from PPTX chart XML", exc_info=True)
 
    return "\n".join(lines)
 
 
def _pptx_image_media_type(content_type: str | None) -> str:
    """Return a valid MIME type for a python-pptx image content_type string."""
    return content_type if content_type else "image/png"